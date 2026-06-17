"""
make_panelist_deck.py
Generates the AI Terminology for Public Health — Panelist Briefing deck.
Run: python3 make_panelist_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1a, 0x28, 0x57)
YELLOW = RGBColor(0xF5, 0xC0, 0x00)
YELLOW_LIGHT = RGBColor(0xFF, 0xF8, 0xD6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
TEXT   = RGBColor(0x0F, 0x17, 0x2A)
MUTED  = RGBColor(0x64, 0x74, 0x8B)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)

OPT_COLORS = [
    RGBColor(0x25, 0x63, 0xEB),  # A — blue
    RGBColor(0x7C, 0x3A, 0xED),  # B — purple
    RGBColor(0x05, 0x96, 0x69),  # C — green
    RGBColor(0xD9, 0x77, 0x06),  # D — amber
]

# ── Helpers ───────────────────────────────────────────────────────────────────
W, H = Inches(13.33), Inches(7.5)   # LAYOUT_WIDE

def rgb(r, g, b): return RGBColor(r, g, b)

def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_w=0):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=14, bold=False, italic=False,
             color=TEXT, align=PP_ALIGN.LEFT, valign=None,
             fill=None, wrap=True, font_face="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if valign:
        from pptx.enum.text import MSO_ANCHOR
        tf.auto_size = None
        mapping = {'middle': MSO_ANCHOR.MIDDLE, 'top': MSO_ANCHOR.TOP, 'bottom': MSO_ANCHOR.BOTTOM}
        tf.vertical_anchor = mapping.get(valign, MSO_ANCHOR.TOP)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_face
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if fill:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = fill
    return txBox

def add_multiline(slide, lines, x, y, w, h, size=13, color=TEXT,
                  bold=False, line_spacing=1.15, font_face="Calibri"):
    """lines: list of (text, bold, color, size) tuples or plain strings"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    for idx, line in enumerate(lines):
        if isinstance(line, str):
            t, b, c, s = line, bold, color, size
        else:
            t = line.get('text', '')
            b = line.get('bold', bold)
            c = line.get('color', color)
            s = line.get('size', size)

        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # line spacing
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = t
        run.font.name = font_face
        run.font.size = _Pt(s)
        run.font.bold = b
        run.font.color.rgb = c
    return txBox

def add_option_card(slide, letter, text, source, color, x, y, w=5.8, h=1.55):
    """Renders a definition option card with colored left bar, letter, text, source."""
    # Card background
    add_rect(slide, x, y, w, h, LIGHT_BG)
    # Left color bar
    add_rect(slide, x, y, 0.07, h, color)
    # Letter badge
    add_rect(slide, x+0.12, y+0.12, 0.32, 0.32, color)
    add_text(slide, letter, x+0.12, y+0.12, 0.32, 0.32,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign='middle')
    # Option text
    add_text(slide, text, x+0.52, y+0.1, w-0.62, h-0.32,
             size=12, color=TEXT, wrap=True)
    # Source label
    if source:
        add_text(slide, source, x+0.52, y+h-0.28, w-0.62, 0.25,
                 size=11, italic=True, color=MUTED)

def add_scenario_box(slide, scenario_text, x, y, w, h):
    """Navy scenario callout box."""
    add_rect(slide, x, y, w, h, NAVY)
    # Yellow left bar
    add_rect(slide, x, y, 0.06, h, YELLOW)
    # Label
    add_text(slide, 'SCENARIO', x+0.15, y+0.12, w-0.2, 0.22,
             size=10, bold=True, color=YELLOW, align=PP_ALIGN.LEFT)
    # Text
    add_text(slide, scenario_text, x+0.15, y+0.38, w-0.2, h-0.48,
             size=12, color=WHITE, wrap=True)

def add_navy_header(slide, term_name=None, right_label=None):
    """Standard navy header bar for content slides."""
    add_rect(slide, 0, 0, 13.33, 0.95, NAVY)
    if term_name:
        # Yellow pill for term
        add_rect(slide, 0.35, 0.18, 0.08, 0.58, YELLOW)  # thin yellow left accent
        add_text(slide, term_name.upper(), 0.52, 0.19, 3.5, 0.55,
                 size=9, bold=True, color=YELLOW, valign='middle')
    if right_label:
        add_text(slide, right_label, 0, 0.1, 13.0, 0.75,
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, valign='middle')
    elif term_name:
        # No right label, just spacing
        pass

def add_content_header(slide, title, term_name=None):
    """Navy header with title + optional yellow term pill."""
    add_rect(slide, 0, 0, 13.33, 0.9, NAVY)
    if term_name:
        # Yellow pill
        add_rect(slide, 0.35, 0.2, 2.4, 0.5, YELLOW)
        add_text(slide, term_name.upper(), 0.35, 0.2, 2.4, 0.5,
                 size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER, valign='middle')
        add_text(slide, title, 2.9, 0.05, 10.0, 0.8,
                 size=18, bold=True, color=WHITE, valign='middle')
    else:
        add_text(slide, title, 0.4, 0.05, 12.5, 0.8,
                 size=20, bold=True, color=WHITE, valign='middle')

# ── Content data ──────────────────────────────────────────────────────────────
# Kept in sync with webinar_app/public/terms.js
TERMS = {
    'hitl': {
        'name': 'Human in the Loop',
        'icon': '👤',
        'concept': 'A system in which a human provides oversight, validation, or correction of an AI system\'s outputs at defined points, retaining decision authority.',
        'ambiguity': 'The term assumes a qualified specialist is available to review AI outputs. In many public and community health settings, the person physically present may lack the expertise the term implies.',
        'question': 'Is any human presence enough to call it "human in the loop", or does the human need the expertise to meaningfully override the AI?',
        'formatA_prompt': 'Four definitions of Human in the Loop. Select the one closest to how YOU understand the term.',
        'options': [
            ('A', 'A risk-mitigation requirement under which a human must be involved in or accountable for AI-driven decisions, including in settings where a specialist who would otherwise make the decision is absent.', 'Public health implementer lens'),
            ('B', 'An AI deployment pattern in which a clinician reviews and either accepts, modifies, or rejects each AI recommendation before it affects patient care.', 'Clinical user lens'),
            ('C', 'A workflow design choice in which AI handles routine cases and escalates uncertain or high-risk cases to a human reviewer, optimising for the limited time a human can spend.', 'Developer / implementer lens'),
            ('D', 'A system in which a person retains decision authority through oversight, validation, or correction of an AI system\'s outputs at one or more defined points in its operation.', 'Regulatory framing'),
        ],
        'scenario': 'Your organization supports an AI tool that reads chest X-rays for diagnosis of TB. The tool is deployed in a district hospital with no radiologist on site. A trained nurse runs the X-ray, the AI flags the result as "high probability TB" and the patient is started on treatment based on the AI\'s output and the nurse\'s clinical assessment.',
        'scenario_q': 'Is this "Human in the Loop"?',
        'scenario_opts': [
            ('A', 'Yes. The nurse is the human, and she is in the loop.'),
            ('B', 'No. The human should be a specialist with the expertise to validate the AI\'s outputs, which the nurse does not have.'),
            ('C', 'It depends on what the term is being used to certify.'),
        ],
        'lightning_q': 'Is Human in the Loop primarily a safety mechanism or a workforce-substitution mechanism?',
        'lightning_opts': [('A', 'A safety mechanism'), ('B', 'A workforce-substitution mechanism')],
        'discussion': [
            'If we accept that the nurse is the human in the loop, is the term doing any work? What is it actually certifying?',
            'If we say it\'s not HITL, are we saying this deployment shouldn\'t happen? Because in many districts, the alternative is no diagnostic at all.',
            'Should the term distinguish between "human supervising AI" (clinician oversight) and "human enabled by AI" (frontline worker using AI to operate beyond their scope)?',
        ],
    },
    'perf': {
        'name': 'Performance',
        'icon': '📊',
        'concept': 'Measures of how well an AI tool performs on a given task, typically derived from model evaluation on a test dataset.',
        'ambiguity': 'An AI model\'s accuracy against a benchmark and a health program\'s impact on a population are not the same claim. What metrics are needed to support evidence across algorithmic accuracy and public health utility + impact?',
        'question': 'Performance according to whom, measured how, and sufficient for what decision?',
        'formatA_prompt': 'Here are four ways to think about performance for an AI tool in public health. Select the one closest to how YOU understand the term.',
        'options': [
            ('A', 'Measures of clinical performance: how well the tool supports correct clinical decisions in real practice, including agreement with expert reviewers.', 'Clinical evaluation lens'),
            ('B', 'Measures of model accuracy on a test dataset, typically including precision, recall, F1 scores, and confusion matrix outputs.', 'ML / developer lens'),
            ('C', 'Measures relevant to procurement and sustainability: uptime, integration cost, frontline user adoption, total cost per case and operational reliability.', 'Health system / Ministry decision-maker lens'),
            ('D', 'Measures of population-level impact: change in detection rates, time-to-treatment, quality of treatment and downstream health outcomes attributable to the tool.', 'Program evaluation lens'),
        ],
        'scenario': 'Your organization has implemented a maternal health WhatsApp chatbot deployed in partnership with the Ministry of Health. After 18 months, you report 92% accuracy on intent recognition, 88% user satisfaction in a 500-person sample, and 1.2 million messages handled. The Ministry is deciding whether to absorb the chatbot into its routine program budget.',
        'scenario_q': 'Are these the right performance metrics for the Ministry\'s decision?',
        'scenario_opts': [
            ('A', 'Yes. These are standard, widely accepted measures.'),
            ('B', 'No. They don\'t tell the Ministry what it needs to know.'),
            ('C', 'Partly. They\'re necessary but not sufficient.'),
        ],
        'lightning_q': 'When you hear "this AI tool achieved 95% accuracy", what should the first follow-up question be?',
        'lightning_opts': [('A', 'On what dataset?'), ('B', 'What changed in program outcomes?')],
        'discussion': [
            'What\'s missing? (equity of access across language groups, behaviour change outcomes, harm cases, cost per enrolled user, sustainability after donor exit.)',
            'If you were advising the Ministry, what is the ONE additional metric you would require before sign-off?',
            'Whose job should it be to produce the metrics that matter? The product developer, the implementing partner, or the Ministry of Health?',
        ],
    },
    'hallu': {
        'name': 'Hallucinations',
        'icon': '⚠️',
        'concept': 'Outputs from a generative AI model that are factually incorrect, fabricated, or context-inappropriate, presented with the same confidence as correct answers.',
        'ambiguity': 'In high-resource settings, a wrong AI output is a patient safety event with a clear responsible party. In public health, frontline workers may lack the reference materials to detect errors, and accountability is rarely clear.',
        'question': 'What is the right comparison group: ideal performance, average clinician performance, or no guidance at all?',
        'formatA_prompt': 'Four ways to define Hallucinations. Select the one closest to how YOU use the term.',
        'options': [
            ('A', 'A safety event in which a generative AI tool, deployed in a clinical or health-information context, gives advice that is plausibly worded but medically wrong, with potential to cause harm.', 'Clinical safety lens'),
            ('B', 'A failure mode of the model-context fit: an output that may be technically accurate against the model\'s training data but is wrong in the specific context where the user is reading it (wrong country, wrong guideline version, wrong population).', 'Deployment / context lens'),
            ('C', 'A trust failure: any output that a user reasonably acts upon, but which is not traceable to a verifiable source, regardless of whether the underlying fact happens to be correct.', 'Public health implementer lens'),
            ('D', 'An output from a generative model that is factually incorrect, fabricated, or unsupported by its training data, produced with a tone of confidence indistinguishable from a correct answer.', 'ML / technical lens'),
        ],
        'scenario': 'A generative AI chatbot is being piloted with Community Health Workers to support quick clinical guidance during home visits. During the pilot, an evaluator finds that for roughly 3% of queries the chatbot gives an answer that contradicts the national clinical guideline. In roughly 1% of queries it contradicts WHO guidance. In all of these cases the answer is presented with the same confidence as the correct answers. The program team is debating whether to continue scaling the pilot.',
        'scenario_q': 'How should the program team respond?',
        'scenario_opts': [
            ('A', 'Continue scaling. 96% guideline-aligned is better than what can be achieved without the tool.'),
            ('B', 'Pause scaling until the rate of guideline-contradicting answers is brought below a defined threshold.'),
            ('C', 'Continue scaling, but with mandatory grounding (RAG against the national guideline) and a visible source citation on every answer.'),
            ('D', 'Reframe the tool as a search aid that surfaces sources, not as a question-answering tool.'),
        ],
        'lightning_q': 'For a Generative AI tool used by Community Health Workers, the most useful mitigation against hallucinations is:',
        'lightning_opts': [('A', 'Better models (fewer hallucinations at the model level)'), ('B', 'Better grounding (answers tied to a verifiable source that can be checked)')],
        'discussion': [
            'Where does responsibility sit when a CHW acts on a hallucinated answer? The CHW, the implementing partner, the vendor, or the Ministry that approved the pilot?',
            'Is "hallucination rate" even the right metric, or does it hide more than it reveals (severity, equity of error distribution, types of question that trigger errors)?',
            'Are we asking generative AI to clear a bar that we never asked older job-aids or printed protocols to clear? What\'s the right comparison group?',
        ],
    },
    'bias': {
        'name': 'Bias',
        'icon': '⚖️',
        'concept': 'Systematic patterns in AI outputs that affect some populations differently from others, often along lines of geography, language, income, or demographic group.',
        'ambiguity': 'Where you locate the source of bias determines who you ask to fix it. A model-bias framing points toward retraining. A data infrastructure framing points toward surveillance investment and governance reform.',
        'question': 'Is this a model problem or a data system problem, and does the answer change what we do about it?',
        'formatA_prompt': 'Four framings of Bias / Fairness in AI for health. Select the one closest to how YOU understand the term.',
        'options': [
            ('A', 'A statistical property of a model: systematic deviation between predictions and ground truth, typically measured by performance differences across demographic subgroups.', 'ML / technical lens'),
            ('B', 'A training data problem: when the data used to build a model under-represents certain populations, the model\'s performance on those populations will be poorer.', 'Data science lens'),
            ('C', 'A system-level problem of who is in the data at all: the populations least visible to digital health systems are the ones most at risk of being further excluded by AI tools built on those systems.', 'Equity / public health lens'),
            ('D', 'A clinical safety issue: when an AI tool produces systematically different recommendations for patients with similar clinical needs but different demographic profiles, leading to inequitable care.', 'Clinical / patient-safety lens'),
        ],
        'scenario': 'A national early-warning system uses AI to flag districts at elevated risk of cholera outbreaks, drawing on routine health facility data, mobile-phone movement patterns, and rainfall data. Six months in, an evaluation finds: the model correctly flagged 9 of 11 outbreaks in urban districts, but missed 4 of 5 outbreaks in nomadic and remote rural areas.',
        'scenario_q': 'Is this a biased model?',
        'scenario_opts': [
            ('A', 'Yes. The model performs unequally across populations.'),
            ('B', 'No. The model is doing the best it can with the data it has. The problem is upstream.'),
            ('C', 'Both. And only naming both is useful.'),
        ],
        'lightning_q': 'When an AI tool performs worse for a marginalised population, the primary problem is:',
        'lightning_opts': [('A', 'In the model'), ('B', 'In the data system the model was built on')],
        'discussion': [
            'If the answer is both, how do we communicate that to a Ministry deciding whether to keep using the tool?',
            'Who is responsible for fixing it? The vendor, the Ministry\'s data team, the donor that funded the underlying surveillance system, or your organization?',
            'Is calling this "model bias" actually unhelpful, because it routes the conversation toward retraining when the real fix is investment in the data infrastructure?',
        ],
    },
    'val': {
        'name': 'Validation',
        'icon': '✅',
        'concept': 'The process of testing an AI model to confirm it performs as intended, typically by evaluating it on data it has not previously seen.',
        'ambiguity': 'Validation is often treated as a one-time pre-deployment step conducted in the development setting. For implementing organisations deploying tools across different populations and contexts, the more pressing question is whether that validation holds locally.',
        'question': 'If a tool was validated somewhere else, how do you know it works here?',
        'formatA_prompt': 'Four definitions of Validation for an AI tool in health. Which one most closely reflects how YOU would apply the term in your work?',
        'options': [
            ('A', 'A continuous process rather than a one-time event: ongoing assessment of whether a deployed model continues to perform as intended as populations, data systems, and care delivery patterns change over time.', 'Post-market surveillance lens'),
            ('B', 'The process of testing a trained AI model on a dataset not used during development, drawn from a different institution, population, or time period, to assess whether its performance holds beyond the original development context.', 'Regulatory / technical lens'),
            ('C', 'The process of assessing whether a model validated internally still performs in a local context on specific populations, facilities and data infrastructure.', 'Public health implementer lens'),
            ('D', 'A one-time pre-deployment requirement: evidence that a model performs adequately on an independent test set before it is approved for clinical or program use.', 'Procurement / approval lens'),
        ],
        'scenario': 'Your organisation is supporting the rollout of an AI clinical scribe across a network of public health facilities. The product documentation states the tool has been validated with high accuracy on clinical consultation transcripts. Your team establishes that all validation datasets came from private hospital settings in two high-income countries with structured consultation formats. Your facilities run high-volume consultations with patient histories that include conditions underrepresented in the training data. Clinicians are required to review and sign off each AI-generated note before it enters the patient record.',
        'scenario_q': 'Is this tool validated for this deployment?',
        'scenario_opts': [
            ('A', 'Yes — the vendor met a recognised standard for validation. What was done is what the term requires.'),
            ('B', 'No — a separate local validation study would be needed before this deployment can be trusted.'),
            ('C', 'The term is being used to mean two different things, and neither party is wrong on their own terms.'),
        ],
        'lightning_q': 'When an AI tool is described as validated, the responsibility for ensuring this applies to the deployment context sits with:',
        'lightning_opts': [
            ('A', 'The product or software vendor, to demonstrate the validation is applicable to the local deployment context.'),
            ('B', 'The implementing organisation, to conduct or commission local validation before deployment.'),
        ],
        'discussion': [
            'If the clinician signs off each note, does that substitute for local validation, or does it transfer the risk to the clinician?',
            'The vendor says validated. The procurement team read that as applicable to this context. Is that a miscommunication, or the predictable result of a term that doesn\'t require context to be specified?',
            'What would a minimum viable local validation study look like for an AI scribe in this setting, and who should fund it?',
        ],
    },
    'gov': {
        'name': 'AI Governance',
        'icon': '🏛️',
        'concept': 'The frameworks, processes, regulations, and accountability structures that guide how AI systems are built, deployed, monitored, and controlled.',
        'ambiguity': 'Regulatory frameworks focus on model risk management. Public health practitioners need to govern the underlying data pipelines, digital infrastructure ownership, and procurement terms that determine who controls AI tools over time.',
        'question': 'Can you meaningfully govern an AI tool without governing the infrastructure it runs on?',
        'formatA_prompt': 'Four definitions of AI Governance in the context of public health. Which one best matches how YOU think it should apply?',
        'options': [
            ('A', 'The internal processes by which an organisation managing or developing an AI system documents its decisions, manages its risks, and ensures accountability for the system\'s behaviour.', 'AI/ML governance lens'),
            ('B', 'The set of laws, regulations, and oversight bodies that authorise, constrain, and audit the deployment of AI in health, including data protection, model approval, and post-market surveillance.', 'Regulatory lens'),
            ('C', 'Government ownership and public stewardship of the foundational digital architecture that determines what AI can be built on, who controls it, and who benefits from it over time.', 'Data sovereignty lens'),
            ('D', 'The framework of policies, ethical principles, transparency requirements, and consent mechanisms that ensure AI in health respects rights, prevents harm, and operates within publicly accountable boundaries.', 'Ethics lens'),
        ],
        'scenario': 'Your organisation supports a national AI program for maternal health in a country. Over three years, the program generates a large dataset of labelled clinical interactions, outcome data, and patient records. The vendor that built and hosts the AI tool uses this data to train an improved global model, licensed commercially to health systems in other countries. The data sharing agreement signed at the start of the program permitted the vendor to use anonymised data for model improvement, but made no distinction between internal improvement and commercial application.',
        'scenario_q': 'Is this an AI governance failure?',
        'scenario_opts': [
            ('A', 'Yes — the vendor exploited a data sharing clause that the Ministry did not intend to include commercial use.'),
            ('B', 'Yes — but the failure was upstream, in the absence of a national framework that could have classified training data as a sovereign asset.'),
            ('C', 'No — the data sharing agreement was signed and complied with. This is a contracting failure, not a governance one.'),
        ],
        'lightning_q': 'When we say AI governance, the primary subject should be:',
        'lightning_opts': [('A', 'The AI model and its behaviour.'), ('B', 'The data and digital infrastructure the AI runs on.')],
        'discussion': [
            'This issue only exists because the data was used to train a model that creates commercial value elsewhere. Does AI governance as a term cover this? If not, what would?',
            'What is the difference between data protection (a legal concept most countries have frameworks for) and data sovereignty (a governance concept most countries do not)?',
            'What role can implementing partners play in helping Ministries negotiate data terms that protect sovereign interests — and is that role currently funded or expected?',
        ],
    },
}

TERM_ORDER = ['hitl', 'perf', 'hallu', 'bias', 'val', 'gov']

# ── Build deck ────────────────────────────────────────────────────────────────
if __name__ == '__main__':
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    # ── SLIDE 1: Cover ────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY

    add_rect(s, 0, 0, 0.45, 7.5, YELLOW)
    add_rect(s, 0.75, 1.6, 3.2, 0.42, YELLOW)
    add_text(s, 'PANELIST BRIEFING', 0.75, 1.6, 3.2, 0.42,
             size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER, valign='middle')
    add_text(s, 'AI Terminology', 0.75, 2.1, 11.0, 1.1,
             size=52, bold=True, color=WHITE, font_face='Calibri')
    add_text(s, 'for Public Health', 0.75, 3.1, 11.0, 0.95,
             size=48, bold=True, color=YELLOW, font_face='Calibri')
    add_rect(s, 0.75, 4.15, 10.5, 0.04, RGBColor(0x40, 0x55, 0x8A))
    add_text(s, '17 June 2026  |  09:00 EST  |  15:00 CAT  |  Online',
             0.75, 4.28, 10.5, 0.45, size=14, color=RGBColor(0xAA, 0xBB, 0xD6))
    add_text(s, 'Led by Daniel Futerman, Project Director CODA, Vital Strategies',
             0.75, 4.85, 8.0, 0.4, size=13, bold=True, color=WHITE)
    add_text(s, 'With invited panelists', 0.75, 5.3, 6.0, 0.35,
             size=12, color=RGBColor(0xAA, 0xBB, 0xD6))

    # ── SLIDE 2: About the session ────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_content_header(s, 'About This Session')

    add_rect(s, 0.5, 1.05, 12.3, 1.0, YELLOW_LIGHT)
    add_rect(s, 0.5, 1.05, 0.06, 1.0, YELLOW)
    add_text(s, 'Does "human in the loop" mean the same thing to a product developer, '
                'an M&E specialist, and a community health worker?',
             0.7, 1.1, 11.9, 0.9, size=14, italic=True, color=NAVY)

    add_text(s, 'Session aims', 0.5, 2.2, 4.0, 0.35, size=13, bold=True, color=NAVY)
    aims = [
        'Surface how common AI terms — built for general technology contexts — get contested when applied to public health practice, especially in resource-limited settings.',
        'Stress-test working definitions against real-world public health scenarios: where do they hold, where do they break?',
        'Move from abstract terminology debates to practical questions: what does this term certify, who is responsible for what it claims, and what decisions does it inform?',
        'Create a live, interactive experience where audience votes and panelist reflections build understanding together.',
    ]
    for i, aim in enumerate(aims):
        add_rect(s, 0.5, 2.65+i*0.98, 0.06, 0.78, OPT_COLORS[i])
        add_text(s, aim, 0.7, 2.67+i*0.98, 11.8, 0.82, size=14, color=TEXT)

    # ── SLIDE 3: Session flow ─────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_content_header(s, 'Session Flow')

    steps = [
        ('1', 'Welcome', NAVY),
        ('2', 'Human in the\nLoop (fixed)', OPT_COLORS[0]),
        ('3', 'Audience votes:\nwhich topics next?', OPT_COLORS[2]),
        ('4', 'Term 2\n(top-voted)', OPT_COLORS[0]),
        ('5', 'Term 3\n(2nd-voted)', OPT_COLORS[0]),
        ('6', 'Optional Term 4,\nif time allows', OPT_COLORS[1]),
        ('7', 'Close', OPT_COLORS[3]),
    ]
    box_w, gap = 1.65, 0.15
    total = len(steps) * box_w + (len(steps)-1) * gap
    start_x = (13.33 - total) / 2

    for i, (num, label, col) in enumerate(steps):
        bx = start_x + i * (box_w + gap)
        if i < len(steps)-1:
            add_rect(s, bx+box_w, 3.0, gap, 0.04, RGBColor(0xCB, 0xD5, 0xE1))
        add_rect(s, bx, 1.1, box_w, 3.7, col)
        add_rect(s, bx+0.6, 1.2, 0.45, 0.45, WHITE)
        add_text(s, num, bx+0.6, 1.2, 0.45, 0.45,
                 size=14, bold=True, color=col, align=PP_ALIGN.CENTER, valign='middle')
        add_text(s, label, bx+0.08, 1.78, box_w-0.16, 2.9,
                 size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign='top')

    add_text(s, 'Every term — fixed or audience-chosen — runs the same five-part format:\n'
                'introduction → definition vote → scenario → lightning vote → panelist discussion & reflections.',
             0.8, 5.1, 11.7, 0.7, size=14, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, 'After Human in the Loop, the audience votes on which of the five remaining terms to explore. '
                'The two highest-voted run automatically; a third runs only if the facilitator judges there\'s time.',
             0.8, 5.85, 11.7, 0.7, size=13, color=MUTED, align=PP_ALIGN.CENTER)

    # ── SLIDE 4: Panelist role ────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_content_header(s, 'Your Role as a Panelist')

    add_text(s, 'What you\'re asked to do', 0.5, 1.1, 5.8, 0.4, size=14, bold=True, color=NAVY)
    panelist_points = [
        'After each term\'s three exercises (definition vote, scenario, lightning vote), you\'ll be invited to share 2-3 minutes of reflection.',
        'React to how the audience voted. What does the split tell you? What would you add or challenge?',
        'Bring your own professional context. You\'re not expected to agree with each other or with the audience.',
        'Keep it grounded: connect your reflections to concrete cases, policy decisions, or implementation realities.',
    ]
    for i, pt in enumerate(panelist_points):
        add_rect(s, 0.5, 1.65+i*1.1, 0.06, 0.85, OPT_COLORS[i])
        add_text(s, pt, 0.7, 1.67+i*1.1, 5.6, 0.88, size=14, color=TEXT)

    add_rect(s, 7.0, 1.0, 5.8, 5.8, LIGHT_BG)
    add_rect(s, 7.0, 1.0, 0.06, 5.8, NAVY)
    add_text(s, 'Participant engagement', 7.2, 1.12, 5.4, 0.38, size=13, bold=True, color=NAVY)
    add_text(s, 'During your discussion segment, participants will:', 7.2, 1.6, 5.3, 0.4,
             size=12, color=MUTED)
    engage_points = [
        ('💬', 'Type questions and comments in the webinar chat'),
        ('✋', 'Raise their hand to speak directly'),
        ('📊', 'See live results from the vote on the shared screen'),
    ]
    for i, (icon, pt) in enumerate(engage_points):
        add_text(s, icon, 7.2, 2.15+i*0.95, 0.5, 0.55, size=18, align=PP_ALIGN.CENTER)
        add_text(s, pt, 7.8, 2.2+i*0.95, 4.6, 0.5, size=14, color=TEXT)

    add_rect(s, 7.1, 5.3, 5.6, 1.15, YELLOW_LIGHT)
    add_rect(s, 7.1, 5.3, 0.06, 1.15, YELLOW)
    add_text(s, 'You don\'t need to review the vote results in advance. The live display during Discussion & Reflections gives you all the context you need.',
             7.3, 5.35, 5.3, 1.05, size=13, italic=True, color=NAVY)

    # ── Term slide helpers ────────────────────────────────────────────────────
    def add_divider(term_key, fixed=False):
        td = TERMS[term_key]
        s = prs.slides.add_slide(blank_layout)
        s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
        add_rect(s, 0, 0, 0.45, 7.5, YELLOW)
        tag = 'TERM 1 — FIXED' if fixed else 'POSSIBLE TOPIC'
        add_rect(s, 0.75, 1.5, 3.5, 0.42, YELLOW)
        add_text(s, tag, 0.75, 1.5, 3.5, 0.42,
                 size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER, valign='middle')
        add_text(s, td['icon'], 9.5, 1.3, 2.5, 2.0, size=72, align=PP_ALIGN.CENTER)
        add_text(s, td['name'], 0.75, 2.05, 8.5, 1.4,
                 size=44, bold=True, color=WHITE, font_face='Calibri')
        add_text(s, td['concept'], 0.75, 3.55, 8.2, 1.3,
                 size=15, color=RGBColor(0xAA, 0xBB, 0xD6), italic=True)
        add_text(s, f'Key question: "{td["question"]}"',
                 0.75, 5.0, 11.5, 1.0, size=15, bold=True, italic=True, color=YELLOW)

    def add_definitions_slide(term_key):
        td = TERMS[term_key]
        s = prs.slides.add_slide(blank_layout)
        s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
        add_content_header(s, 'Exercise 1 — How Do You Define It?', td['name'])
        add_rect(s, 0.4, 1.05, 12.5, 0.78, YELLOW_LIGHT)
        add_rect(s, 0.4, 1.05, 0.06, 0.78, YELLOW)
        add_text(s, td['concept'], 0.6, 1.1, 12.1, 0.68, size=13, italic=True, color=NAVY)
        add_text(s, td['formatA_prompt'], 0.4, 1.94, 12.5, 0.38, size=13, color=MUTED)
        positions = [(0.4, 2.32), (6.85, 2.32), (0.4, 4.58), (6.85, 4.58)]
        for (ox, oy), (letter, text, source) in zip(positions, td['options']):
            col = OPT_COLORS[ord(letter)-ord('A')]
            add_option_card(s, letter, text, source, col, ox, oy, w=6.2, h=2.05)

    def add_scenario_lightning_slide(term_key):
        td = TERMS[term_key]
        s = prs.slides.add_slide(blank_layout)
        s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
        add_content_header(s, 'Exercises 2 & 3 — Scenario and Lightning Vote', td['name'])

        add_rect(s, 0.4, 1.05, 0.06, 5.4, OPT_COLORS[0])
        add_text(s, 'EXERCISE 2 — SCENARIO', 0.6, 1.05, 5.8, 0.32,
                 size=11, bold=True, color=OPT_COLORS[0], align=PP_ALIGN.LEFT)
        add_scenario_box(s, td['scenario'], 0.4, 1.44, 6.0, 2.2)
        add_text(s, td['scenario_q'], 0.4, 3.72, 6.0, 0.42, size=15, bold=True, color=NAVY)
        for i, (letter, text) in enumerate(td['scenario_opts']):
            col = OPT_COLORS[i]
            oy = 4.22 + i * 0.78
            add_rect(s, 0.4, oy, 0.06, 0.62, col)
            add_text(s, letter, 0.52, oy, 0.28, 0.62, size=13, bold=True, color=col, valign='middle')
            add_text(s, text, 0.84, oy, 5.5, 0.62, size=13, color=TEXT, valign='middle')

        add_rect(s, 6.7, 1.05, 0.04, 5.9, BORDER)

        add_rect(s, 6.85, 1.05, 0.06, 5.4, OPT_COLORS[3])
        add_text(s, 'EXERCISE 3 — LIGHTNING VOTE', 7.05, 1.05, 5.8, 0.32,
                 size=11, bold=True, color=OPT_COLORS[3])
        add_text(s, td['lightning_q'], 6.85, 1.44, 6.05, 1.0,
                 size=16, bold=True, color=NAVY)
        lw = 2.8
        for i, (letter, text) in enumerate(td['lightning_opts']):
            col = OPT_COLORS[i]
            lx = 6.85 + i * (lw + 0.25)
            add_rect(s, lx, 2.58, lw, 2.7, LIGHT_BG)
            add_rect(s, lx, 2.58, lw, 0.07, col)
            add_text(s, letter, lx, 2.72, lw, 0.55,
                     size=26, bold=True, color=col, align=PP_ALIGN.CENTER)
            add_text(s, text, lx+0.12, 3.3, lw-0.24, 1.5,
                     size=14, color=TEXT, align=PP_ALIGN.CENTER)

        add_rect(s, 6.85, 5.45, 6.05, 1.75, LIGHT_BG)
        add_rect(s, 6.85, 5.45, 0.06, 1.75, NAVY)
        add_text(s, 'DISCUSSION PROMPTS FOR PANELISTS', 7.05, 5.5, 5.7, 0.32,
                 size=10, bold=True, color=NAVY)
        for i, dp in enumerate(td['discussion']):
            short = dp[:120] + ('…' if len(dp) > 120 else '')
            add_text(s, f'• {short}', 7.05, 5.87+i*0.44, 5.7, 0.4, size=11, color=MUTED)
            if i >= 2: break

    # ── BUILD ALL SLIDES ──────────────────────────────────────────────────────
    add_divider('hitl', fixed=True)
    add_definitions_slide('hitl')
    add_scenario_lightning_slide('hitl')

    for tk in TERM_ORDER[1:]:
        add_divider(tk, fixed=False)
        add_definitions_slide(tk)
        add_scenario_lightning_slide(tk)

    # --- Prep guide ---
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_content_header(s, 'Preparing for the Session')

    prep_items = [
        ('Skim each term section', 'Read the concept, the ambiguity framing, and the key question for each of the six terms. You don\'t need to prepare formal remarks.'),
        ('Think about your own context', 'For each term, consider: have you encountered this in your work? Where has the definition mattered? Where has it caused confusion or friction?'),
        ('Scenarios are the richest territory', 'The scenario exercises are where definitions get stress-tested. Think about analogous cases from your own experience you could reference.'),
        ('You don\'t need to agree', 'Disagreement between panelists on the same question is exactly what makes this session valuable for the audience.'),
        ('Lightning votes are conversation starters', 'The binary vote is intentionally provocative. The point is not the answer but the question it raises.'),
    ]
    for i, (title, desc) in enumerate(prep_items):
        oy = 1.1 + i * 1.2
        add_rect(s, 0.4, oy, 0.06, 1.0, OPT_COLORS[i % 4])
        add_text(s, title, 0.6, oy+0.05, 4.5, 0.38, size=14, bold=True, color=NAVY)
        add_text(s, desc, 0.6, oy+0.45, 12.0, 0.55, size=13, color=MUTED)

    add_rect(s, 0, 7.1, 13.33, 0.4, NAVY)
    add_text(s, 'Questions before the session? Contact Daniel Futerman, Vital Strategies',
             0.4, 7.12, 12.5, 0.3, size=12, color=RGBColor(0xAA, 0xBB, 0xD6))

    # ── Save ──────────────────────────────────────────────────────────────────
    out = '/Users/danielfuterman/code/aiterm/AI_Terminology_Panelist_Briefing.pptx'
    prs.save(out)
    print(f'Saved: {out}')
    print(f'Slides: {len(prs.slides)}')
