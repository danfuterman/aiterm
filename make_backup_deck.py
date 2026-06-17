"""
make_backup_deck.py
Generates the AI Terminology for Public Health — Facilitator Backup Deck.
Use this if the live voting app is unavailable during the session.
Run: python3 make_backup_deck.py
"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from make_panelist_deck import (
    NAVY, YELLOW, YELLOW_LIGHT, WHITE, LIGHT_BG, TEXT, MUTED, BORDER, OPT_COLORS,
    add_rect, add_text, add_scenario_box, add_content_header,
    TERMS, TERM_ORDER,
)

# ── Deck setup ────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_option_card_wide(slide, letter, text, color, x, y, w=12.5, h=1.28):
    """Full-width option card — no source label (audience-facing)."""
    add_rect(slide, x, y, w, h, LIGHT_BG)
    add_rect(slide, x, y, 0.07, h, color)
    badge_y = y + (h - 0.40) / 2
    add_rect(slide, x+0.14, badge_y, 0.40, 0.40, color)
    add_text(slide, letter, x+0.14, badge_y, 0.40, 0.40,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign='middle')
    add_text(slide, text, x+0.66, y+0.1, w-0.76, h-0.2, size=16, color=TEXT, wrap=True)


# ── SLIDE 1: Cover ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
add_rect(s, 0, 0, 0.45, 7.5, YELLOW)
add_rect(s, 0.75, 1.6, 5.2, 0.44, YELLOW)
add_text(s, 'FACILITATOR BACKUP DECK', 0.75, 1.6, 5.2, 0.44,
         size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER, valign='middle')
add_text(s, 'AI Terminology', 0.75, 2.18, 11.0, 1.1,
         size=52, bold=True, color=WHITE, font_face='Calibri')
add_text(s, 'for Public Health', 0.75, 3.18, 11.0, 0.95,
         size=48, bold=True, color=YELLOW, font_face='Calibri')
add_rect(s, 0.75, 4.24, 10.5, 0.04, RGBColor(0x40, 0x55, 0x8A))
add_text(s, '17 June 2026  |  09:00 EST  |  15:00 CAT  |  Online',
         0.75, 4.38, 10.5, 0.45, size=14, color=RGBColor(0xAA, 0xBB, 0xD6))
add_text(s, 'Use this deck if the live voting app is unavailable.',
         0.75, 5.05, 10.5, 0.45, size=14, italic=True, color=RGBColor(0xAA, 0xBB, 0xD6))

# ── SLIDE 2: Session flow ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
add_content_header(s, 'Session Flow')
steps = [
    ('1', 'Welcome', NAVY),
    ('2', 'Human in the\nLoop (fixed)', OPT_COLORS[0]),
    ('3', 'Audience votes:\nwhich topics next?', OPT_COLORS[2]),
    ('4', 'Term 2\n(top-voted)', OPT_COLORS[0]),
    ('5', 'Term 3\n(2nd-voted)', OPT_COLORS[0]),
    ('6', 'Optional Term 4,\nif time allows', OPT_COLORS[1]),
    ('7', 'Close', OPT_COLORS[3]),
]
box_w, gap = 1.65, 0.15
start_x = (13.33 - (len(steps) * box_w + (len(steps)-1) * gap)) / 2
for i, (num, label, col) in enumerate(steps):
    bx = start_x + i * (box_w + gap)
    if i < len(steps)-1:
        add_rect(s, bx+box_w, 3.0, gap, 0.04, RGBColor(0xCB, 0xD5, 0xE1))
    add_rect(s, bx, 1.1, box_w, 3.7, col)
    add_rect(s, bx+0.6, 1.2, 0.45, 0.45, WHITE)
    add_text(s, num, bx+0.6, 1.2, 0.45, 0.45,
             size=14, bold=True, color=col, align=PP_ALIGN.CENTER, valign='middle')
    add_text(s, label, bx+0.08, 1.78, box_w-0.16, 2.9,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign='top')
add_text(s, 'Each term: introduction → definition vote → scenario → lightning vote → discussion.',
         0.8, 5.1, 11.7, 0.5, size=15, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# ── SLIDE 3: Topic vote (shortlist) ──────────────────────────────────────────
s = prs.slides.add_slide(blank)
s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
add_content_header(s, 'Vote: Which topics should we cover next?')
add_text(s, 'After Human in the Loop, indicate your top 2 topic choices.',
         0.5, 1.05, 12.3, 0.4, size=16, italic=True, color=MUTED)
shortlist = [k for k in TERM_ORDER if k != 'hitl']
for i, tk in enumerate(shortlist):
    td = TERMS[tk]
    col = OPT_COLORS[i % 4]
    oy = 1.6 + i * 1.12
    add_rect(s, 0.5, oy, 12.3, 1.0, LIGHT_BG)
    add_rect(s, 0.5, oy, 0.07, 1.0, col)
    add_rect(s, 0.68, oy+0.12, 0.48, 0.48, col)
    add_text(s, 'ABCDE'[i], 0.68, oy+0.12, 0.48, 0.48,
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign='middle')
    add_text(s, td['name'], 1.3, oy+0.1, 4.5, 0.38,
             size=17, bold=True, color=NAVY)
    add_text(s, td['concept'][:120] + ('…' if len(td['concept']) > 120 else ''),
             1.3, oy+0.54, 11.1, 0.38, size=13, italic=True, color=MUTED)

# ── Per-term slides ───────────────────────────────────────────────────────────
for tk in TERM_ORDER:
    td   = TERMS[tk]
    tag  = 'TERM 1 — FIXED' if tk == 'hitl' else 'POSSIBLE TOPIC'

    # Slide A: Term divider (navy, large)
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
    add_rect(s, 0, 0, 0.45, 7.5, YELLOW)
    add_rect(s, 0.75, 1.5, 3.5, 0.44, YELLOW)
    add_text(s, tag, 0.75, 1.5, 3.5, 0.44,
             size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER, valign='middle')
    add_text(s, td['icon'], 9.5, 1.3, 2.5, 2.0, size=72, align=PP_ALIGN.CENTER)
    add_text(s, td['name'], 0.75, 2.1, 8.5, 1.35,
             size=44, bold=True, color=WHITE, font_face='Calibri')
    add_text(s, td['concept'], 0.75, 3.58, 8.5, 1.3,
             size=16, italic=True, color=RGBColor(0xAA, 0xBB, 0xD6))
    add_text(s, f'Key question: "{td["question"]}"',
             0.75, 5.1, 11.5, 1.1, size=16, bold=True, italic=True, color=YELLOW)

    # Slide B: Definition vote — full-width vertical list
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_content_header(s, 'How Do You Define It?', td['name'])
    add_rect(s, 0.4, 1.05, 12.5, 0.72, YELLOW_LIGHT)
    add_rect(s, 0.4, 1.05, 0.06, 0.72, YELLOW)
    add_text(s, td['formatA_prompt'], 0.6, 1.1, 12.1, 0.62, size=14, italic=True, color=NAVY)

    n_opts  = len(td['options'])
    opt_top = 1.88
    opt_h   = (7.5 - opt_top - 0.15) / n_opts
    for idx, (letter, text, source) in enumerate(td['options']):
        col = OPT_COLORS[idx]
        add_option_card_wide(s, letter, text, col, 0.4, opt_top + idx * opt_h,
                             w=12.5, h=opt_h - 0.08)

    # Slide C: Scenario
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_content_header(s, 'Scenario', td['name'])
    sc_box_h = 2.25
    add_scenario_box(s, td['scenario'], 0.4, 1.08, 12.5, sc_box_h)
    q_top = 1.08 + sc_box_h + 0.14
    add_text(s, td['scenario_q'], 0.4, q_top, 12.5, 0.5, size=18, bold=True, color=NAVY)
    sc_opts   = td['scenario_opts']
    n_sc      = len(sc_opts)
    sc_top    = q_top + 0.56
    sc_opt_h  = (7.5 - sc_top - 0.15) / n_sc
    for idx, (letter, text) in enumerate(sc_opts):
        col = OPT_COLORS[idx]
        add_option_card_wide(s, letter, text, col, 0.4, sc_top + idx * sc_opt_h,
                             w=12.5, h=sc_opt_h - 0.1)

    # Slide D: Lightning vote — two large cards
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_content_header(s, 'Lightning Vote', td['name'])
    add_text(s, td['lightning_q'], 0.5, 1.08, 12.3, 1.2,
             size=22, bold=True, color=NAVY)
    lw      = (13.33 - 0.8 - 0.35) / 2
    lcard_h = 7.5 - 2.55 - 0.2
    for i, (letter, text) in enumerate(td['lightning_opts']):
        col = OPT_COLORS[i]
        lx  = 0.4 + i * (lw + 0.35)
        add_rect(s, lx, 2.45, lw, lcard_h, LIGHT_BG)
        add_rect(s, lx, 2.45, lw, 0.1, col)
        add_text(s, letter, lx, 2.65, lw, 0.85,
                 size=42, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(s, text, lx+0.2, 3.58, lw-0.4, lcard_h-1.22,
                 size=19, color=TEXT, align=PP_ALIGN.CENTER, wrap=True)

    # Slide E: Discussion & Reflections
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_content_header(s, 'Discussion & Reflections', td['name'])
    n_dp  = min(len(td['discussion']), 3)
    dp_h  = (7.5 - 1.2 - 0.3) / n_dp
    for i, dp in enumerate(td['discussion'][:3]):
        col = OPT_COLORS[i]
        oy  = 1.2 + i * dp_h
        add_rect(s, 0.4, oy+0.1, 0.08, dp_h-0.25, col)
        add_text(s, dp, 0.65, oy+0.1, 12.0, dp_h-0.3, size=17, color=TEXT, wrap=True)

# ── Save ──────────────────────────────────────────────────────────────────────
out = '/Users/danielfuterman/code/aiterm/AI_Terminology_Backup_Deck.pptx'
prs.save(out)
print(f'Saved: {out}')
print(f'Slides: {len(prs.slides)}')
